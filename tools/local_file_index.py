#!/usr/bin/env python
"""Local file index for agent-readable desktop search.

The index is local-only by default. It records searchable text in a SQLite FTS
database under tmp/local-file-index/, and does not upload content to EverOS.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import os
import re
import sqlite3
import sys
import time
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_DB = PROJECT_ROOT / "tmp" / "local-file-index" / "files.sqlite3"

DEFAULT_ALLOWED_EXTENSIONS = {
    ".bat",
    ".cfg",
    ".c",
    ".cc",
    ".cpp",
    ".cs",
    ".css",
    ".csv",
    ".docx",
    ".go",
    ".h",
    ".hpp",
    ".htm",
    ".html",
    ".ini",
    ".java",
    ".js",
    ".json",
    ".jsonl",
    ".jsx",
    ".kt",
    ".log",
    ".md",
    ".pdf",
    ".php",
    ".pptx",
    ".ps1",
    ".py",
    ".rb",
    ".rs",
    ".sh",
    ".sql",
    ".toml",
    ".ts",
    ".tsx",
    ".txt",
    ".xml",
    ".xlsx",
    ".yaml",
    ".yml",
}

DEFAULT_SKIP_DIR_NAMES = {
    "$recycle.bin",
    ".cache",
    ".git",
    ".hg",
    ".idea",
    ".mypy_cache",
    ".pytest_cache",
    ".svn",
    ".tox",
    ".venv",
    ".vscode",
    "__pycache__",
    "appdata",
    "node_modules",
    "program files",
    "program files (x86)",
    "programdata",
    "system volume information",
    "venv",
    "windows",
}

DEFAULT_SKIP_FILE_NAMES = {
    ".env",
    ".env.local",
    ".npmrc",
    ".pypirc",
    "id_dsa",
    "id_ecdsa",
    "id_ed25519",
    "id_rsa",
    "known_hosts",
}

DEFAULT_SENSITIVE_EXTENSIONS = {
    ".cer",
    ".crt",
    ".db",
    ".der",
    ".key",
    ".kdbx",
    ".p12",
    ".pem",
    ".pfx",
    ".sqlite",
    ".sqlite3",
}


@dataclass(frozen=True)
class ExtractedText:
    status: str
    text: str
    error: str = ""


def norm_ext(ext: str) -> str:
    ext = ext.strip().lower()
    return ext if ext.startswith(".") else f".{ext}"


def parse_extensions(raw: str | None) -> set[str]:
    if not raw:
        return set(DEFAULT_ALLOWED_EXTENSIONS)
    return {norm_ext(item) for item in raw.split(",") if item.strip()}


def connect(db_path: Path) -> sqlite3.Connection:
    db_path.parent.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(db_path)
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA synchronous=NORMAL")
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS files (
          path TEXT PRIMARY KEY,
          name TEXT NOT NULL,
          suffix TEXT NOT NULL,
          size INTEGER NOT NULL,
          mtime_ns INTEGER NOT NULL,
          sha256 TEXT,
          indexed_at INTEGER NOT NULL,
          status TEXT NOT NULL,
          error TEXT,
          content_chars INTEGER NOT NULL
        )
        """
    )
    conn.execute(
        """
        CREATE VIRTUAL TABLE IF NOT EXISTS file_fts
        USING fts5(path UNINDEXED, name, content, tokenize='unicode61')
        """
    )
    return conn


def is_sensitive_file(path: Path) -> bool:
    lower_name = path.name.lower()
    if lower_name in DEFAULT_SKIP_FILE_NAMES:
        return True
    if lower_name.startswith(".env."):
        return True
    if path.suffix.lower() in DEFAULT_SENSITIVE_EXTENSIONS:
        return True
    return False


def should_skip_dir(path: Path, extra_skip_dirs: set[str]) -> bool:
    name = path.name.lower()
    return name in DEFAULT_SKIP_DIR_NAMES or name in extra_skip_dirs


def common_roots() -> list[Path]:
    home = Path.home()
    candidates = [
        PROJECT_ROOT,
        Path("D:/Docs"),
        Path("D:/Projects"),
        home / "Desktop",
        home / "Documents",
        home / "Downloads",
    ]
    return unique_existing_paths(candidates)


def drive_roots() -> list[Path]:
    roots: list[Path] = []
    for code in range(ord("A"), ord("Z") + 1):
        root = Path(f"{chr(code)}:/")
        if root.exists():
            roots.append(root)
    return roots


def unique_existing_paths(paths: Iterable[Path]) -> list[Path]:
    seen: set[str] = set()
    result: list[Path] = []
    for path in paths:
        try:
            resolved = path.resolve()
        except OSError:
            continue
        key = str(resolved).casefold()
        if resolved.exists() and key not in seen:
            seen.add(key)
            result.append(resolved)
    return result


def resolve_roots(args: argparse.Namespace) -> list[Path]:
    roots = [Path(item) for item in args.root]
    if args.preset == "common":
        roots.extend(common_roots())
    elif args.preset == "all-drives":
        if not args.allow_all_drives:
            raise SystemExit("--preset all-drives requires --allow-all-drives.")
        roots.extend(drive_roots())

    roots = unique_existing_paths(roots)
    if not roots:
        raise SystemExit("No existing roots selected. Pass --root PATH or --preset common.")
    return roots


def iter_files(
    roots: list[Path],
    allowed_extensions: set[str],
    extra_skip_dirs: set[str],
    max_bytes: int,
    limit: int | None = None,
) -> Iterable[Path]:
    yielded = 0
    for root in roots:
        for dirpath, dirnames, filenames in os.walk(root, topdown=True, onerror=lambda err: None):
            current = Path(dirpath)
            dirnames[:] = [
                name for name in dirnames if not should_skip_dir(current / name, extra_skip_dirs)
            ]
            for filename in filenames:
                path = current / filename
                suffix = path.suffix.lower()
                if suffix not in allowed_extensions:
                    continue
                if is_sensitive_file(path):
                    continue
                try:
                    stat = path.stat()
                except OSError:
                    continue
                if stat.st_size > max_bytes:
                    continue
                yielded += 1
                yield path
                if limit and yielded >= limit:
                    return


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest().upper()


def trim_text(text: str, max_chars: int) -> str:
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", " ", text)
    return text[:max_chars]


def read_text_file(path: Path, max_chars: int) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return trim_text(raw.decode(encoding), max_chars)
        except UnicodeDecodeError:
            continue
    return trim_text(raw.decode("utf-8", errors="replace"), max_chars)


def extract_pdf(path: Path, max_chars: int) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    total = 0
    for page in reader.pages:
        text = page.extract_text() or ""
        if not text:
            continue
        parts.append(text)
        total += len(text)
        if total >= max_chars:
            break
    return trim_text("\n".join(parts), max_chars)


def extract_docx(path: Path, max_chars: int) -> str:
    import docx

    document = docx.Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                parts.append(" | ".join(values))
    return trim_text("\n".join(parts), max_chars)


def extract_pptx(path: Path, max_chars: int) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[slide {slide_index}]")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                parts.append(text)
        if sum(len(item) for item in parts) >= max_chars:
            break
    return trim_text("\n".join(parts), max_chars)


def extract_xlsx(path: Path, max_chars: int) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    total = 0
    try:
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            parts.append(f"[sheet {sheet_name}]")
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value not in (None, "")]
                if values:
                    line = "\t".join(values)
                    parts.append(line)
                    total += len(line)
                if total >= max_chars:
                    break
            if total >= max_chars:
                break
    finally:
        workbook.close()
    return trim_text("\n".join(parts), max_chars)


def extract_zip_xml_text(path: Path, max_chars: int) -> str:
    """Fallback for Office XML packages when optional libraries fail."""
    parts: list[str] = []
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            if not name.endswith(".xml"):
                continue
            if not (
                name.startswith("word/")
                or name.startswith("ppt/")
                or name.startswith("xl/sharedStrings")
                or name.startswith("xl/worksheets/")
            ):
                continue
            try:
                root = ElementTree.fromstring(archive.read(name))
            except ElementTree.ParseError:
                continue
            for element in root.iter():
                if element.text and element.text.strip():
                    parts.append(element.text.strip())
            if sum(len(item) for item in parts) >= max_chars:
                break
    return trim_text("\n".join(parts), max_chars)


def extract_text(path: Path, max_chars: int) -> ExtractedText:
    suffix = path.suffix.lower()
    try:
        if suffix in {".pdf"}:
            text = extract_pdf(path, max_chars)
        elif suffix == ".docx":
            try:
                text = extract_docx(path, max_chars)
            except Exception:
                text = extract_zip_xml_text(path, max_chars)
        elif suffix == ".pptx":
            try:
                text = extract_pptx(path, max_chars)
            except Exception:
                text = extract_zip_xml_text(path, max_chars)
        elif suffix == ".xlsx":
            try:
                text = extract_xlsx(path, max_chars)
            except Exception:
                text = extract_zip_xml_text(path, max_chars)
        else:
            text = read_text_file(path, max_chars)
        return ExtractedText(status="indexed", text=text)
    except Exception as exc:
        return ExtractedText(status="error", text="", error=f"{type(exc).__name__}: {exc}")


def upsert_file(conn: sqlite3.Connection, path: Path, extracted: ExtractedText) -> None:
    stat = path.stat()
    sha256 = file_sha256(path)
    content = extracted.text
    payload = (
        str(path.resolve()),
        path.name,
        path.suffix.lower(),
        stat.st_size,
        stat.st_mtime_ns,
        sha256,
        int(time.time()),
        extracted.status,
        extracted.error,
        len(content),
    )
    conn.execute(
        """
        INSERT INTO files(path, name, suffix, size, mtime_ns, sha256, indexed_at, status, error, content_chars)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ON CONFLICT(path) DO UPDATE SET
          name=excluded.name,
          suffix=excluded.suffix,
          size=excluded.size,
          mtime_ns=excluded.mtime_ns,
          sha256=excluded.sha256,
          indexed_at=excluded.indexed_at,
          status=excluded.status,
          error=excluded.error,
          content_chars=excluded.content_chars
        """,
        payload,
    )
    conn.execute("DELETE FROM file_fts WHERE path = ?", (str(path.resolve()),))
    if content:
        conn.execute(
            "INSERT INTO file_fts(path, name, content) VALUES (?, ?, ?)",
            (str(path.resolve()), path.name, content),
        )


def escape_fts_query(query: str) -> str:
    tokens = [token for token in re.split(r"\s+", query.strip()) if token]
    if not tokens:
        raise SystemExit("Search query is empty.")
    escaped = [f'"{token.replace(chr(34), chr(34) + chr(34))}"' for token in tokens]
    return " AND ".join(escaped)


def cmd_scan(args: argparse.Namespace) -> int:
    roots = resolve_roots(args)
    allowed_extensions = parse_extensions(args.extensions)
    extra_skip_dirs = {item.lower() for item in args.exclude_dir}
    max_bytes = int(args.max_mb * 1024 * 1024)

    if args.dry_run:
        count = sum(
            1
            for _ in iter_files(
                roots, allowed_extensions, extra_skip_dirs, max_bytes=max_bytes, limit=args.limit
            )
        )
        print(json.dumps({"dry_run": True, "roots": [str(root) for root in roots], "files": count}, ensure_ascii=False, indent=2))
        return 0

    conn = connect(Path(args.db))
    indexed = 0
    errors = 0
    started = time.time()
    try:
        for path in iter_files(
            roots, allowed_extensions, extra_skip_dirs, max_bytes=max_bytes, limit=args.limit
        ):
            extracted = extract_text(path, args.max_chars)
            if extracted.status == "error":
                errors += 1
            upsert_file(conn, path, extracted)
            indexed += 1
            if indexed % args.commit_every == 0:
                conn.commit()
                print(f"indexed={indexed} errors={errors} current={path}", file=sys.stderr)
        conn.commit()
    finally:
        conn.close()

    print(
        json.dumps(
            {
                "db": str(Path(args.db).resolve()),
                "roots": [str(root) for root in roots],
                "indexed": indexed,
                "errors": errors,
                "elapsed_seconds": round(time.time() - started, 2),
            },
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0


def cmd_search(args: argparse.Namespace) -> int:
    db = Path(args.db)
    if not db.exists():
        raise SystemExit(f"Index not found: {db}. Run scan first.")
    conn = connect(db)
    conn.row_factory = sqlite3.Row
    rows: list[sqlite3.Row] = []
    try:
        fts_query = escape_fts_query(args.query)
        rows = conn.execute(
            """
            SELECT files.path, files.name, files.suffix, files.size, files.status,
                   snippet(file_fts, 2, '[', ']', '...', 12) AS snippet,
                   bm25(file_fts) AS rank
            FROM file_fts
            JOIN files ON files.path = file_fts.path
            WHERE file_fts MATCH ?
            ORDER BY rank
            LIMIT ?
            """,
            (fts_query, args.limit),
        ).fetchall()
    except sqlite3.OperationalError as exc:
        if args.strict:
            raise SystemExit(f"FTS query failed: {exc}") from exc

    if not rows and args.like_fallback:
        needle = f"%{args.query}%"
        rows = conn.execute(
            """
            SELECT files.path, files.name, files.suffix, files.size, files.status,
                   substr(file_fts.content, 1, 240) AS snippet,
                   0 AS rank
            FROM file_fts
            JOIN files ON files.path = file_fts.path
            WHERE file_fts.content LIKE ? OR files.name LIKE ? OR files.path LIKE ?
            LIMIT ?
            """,
            (needle, needle, needle, args.limit),
        ).fetchall()

    if args.format == "json":
        print(json.dumps([dict(row) for row in rows], ensure_ascii=False, indent=2))
    else:
        writer = csv.writer(sys.stdout)
        writer.writerow(["path", "name", "suffix", "size", "status", "snippet"])
        for row in rows:
            writer.writerow([row["path"], row["name"], row["suffix"], row["size"], row["status"], row["snippet"]])
    conn.close()
    return 0


def cmd_show(args: argparse.Namespace) -> int:
    path = Path(args.path)
    if args.path_from_index:
        conn = connect(Path(args.db))
        conn.row_factory = sqlite3.Row
        row = conn.execute("SELECT * FROM files WHERE path = ?", (args.path,)).fetchone()
        conn.close()
        if not row:
            raise SystemExit(f"Path not found in index: {args.path}")
        path = Path(row["path"])

    if args.extract:
        extracted = extract_text(path, args.max_chars)
        print(extracted.text)
        if extracted.error:
            print(extracted.error, file=sys.stderr)
    else:
        print(path.resolve())
    return 0


def cmd_stats(args: argparse.Namespace) -> int:
    db = Path(args.db)
    if not db.exists():
        raise SystemExit(f"Index not found: {db}. Run scan first.")
    conn = connect(db)
    conn.row_factory = sqlite3.Row
    totals = conn.execute(
        "SELECT status, COUNT(*) AS count, SUM(size) AS bytes, SUM(content_chars) AS chars FROM files GROUP BY status"
    ).fetchall()
    suffixes = conn.execute(
        "SELECT suffix, COUNT(*) AS count FROM files GROUP BY suffix ORDER BY count DESC LIMIT 20"
    ).fetchall()
    print(
        json.dumps(
            {
                "db": str(db.resolve()),
                "status": [dict(row) for row in totals],
                "top_suffixes": [dict(row) for row in suffixes],
            },
            ensure_ascii=False,
            indent=2,
        )
    )
    conn.close()
    return 0


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Build and query a local-only desktop file index.")
    parser.add_argument("--db", default=str(DEFAULT_DB), help="SQLite index path.")
    subparsers = parser.add_subparsers(dest="command", required=True)

    scan = subparsers.add_parser("scan", help="Scan selected roots into the local SQLite FTS index.")
    scan.add_argument("--root", action="append", default=[], help="Root directory to scan. Can be repeated.")
    scan.add_argument("--preset", choices=["common", "all-drives"], help="Add a predefined root set.")
    scan.add_argument("--allow-all-drives", action="store_true", help="Required with --preset all-drives.")
    scan.add_argument("--extensions", help="Comma-separated extension allowlist. Default covers text/code/PDF/Office.")
    scan.add_argument("--exclude-dir", action="append", default=[], help="Extra directory name to skip.")
    scan.add_argument("--max-mb", type=float, default=25.0, help="Skip files larger than this size.")
    scan.add_argument("--max-chars", type=int, default=200_000, help="Max extracted characters per file.")
    scan.add_argument("--limit", type=int, help="Stop after N candidate files.")
    scan.add_argument("--commit-every", type=int, default=100)
    scan.add_argument("--dry-run", action="store_true", help="Count candidate files without writing the index.")
    scan.set_defaults(func=cmd_scan)

    search = subparsers.add_parser("search", help="Search indexed files.")
    search.add_argument("query")
    search.add_argument("--limit", type=int, default=20)
    search.add_argument("--format", choices=["csv", "json"], default="csv")
    search.add_argument("--strict", action="store_true", help="Fail instead of falling back when FTS syntax fails.")
    search.add_argument("--like-fallback", action=argparse.BooleanOptionalAction, default=True)
    search.set_defaults(func=cmd_search)

    show = subparsers.add_parser("show", help="Print a path or extract text from a path.")
    show.add_argument("path")
    show.add_argument("--path-from-index", action="store_true")
    show.add_argument("--extract", action="store_true")
    show.add_argument("--max-chars", type=int, default=12000)
    show.set_defaults(func=cmd_show)

    stats = subparsers.add_parser("stats", help="Show index statistics.")
    stats.set_defaults(func=cmd_stats)

    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    return args.func(args)


if __name__ == "__main__":
    raise SystemExit(main())
